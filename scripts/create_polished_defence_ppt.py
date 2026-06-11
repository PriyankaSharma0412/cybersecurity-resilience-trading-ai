from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Priyanka_Dissertation_Defence_Presentation.pptx"
FIG = ROOT / "dissertation_outputs" / "figures"
CH3 = FIG / "chapter3_visuals"
CH4 = FIG / "chapter4_docx_figures"

NAVY = RGBColor(18, 37, 63)
INK = RGBColor(31, 43, 61)
MUTED = RGBColor(89, 104, 124)
BLUE = RGBColor(41, 98, 142)
TEAL = RGBColor(28, 132, 126)
GOLD = RGBColor(187, 132, 32)
RED = RGBColor(178, 76, 76)
BG = RGBColor(247, 249, 252)
PANEL = RGBColor(235, 241, 247)
WHITE = RGBColor(255, 255, 255)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size,
    color=INK,
    bold=False,
    font="Aptos",
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets, x, y, w, h, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.clear()
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.text = f"• {item}"
    return box


def add_rule(slide, x, y, w, color=BLUE):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_tag(slide, text, x, y, w=1.55, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return shape


def add_metric(slide, value, label, x, y, w=1.35, accent=BLUE):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.82))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = accent
    panel.line.width = Pt(1.25)
    add_text(slide, value, x + 0.08, y + 0.08, w - 0.16, 0.28, 20, accent, True)
    add_text(slide, label, x + 0.08, y + 0.42, w - 0.16, 0.18, 10, MUTED, False)


def add_picture_fit(slide, path, x, y, w, h):
    from PIL import Image

    img = Image.open(path)
    iw, ih = img.size
    aspect = iw / ih
    target_aspect = w / h
    if aspect >= target_aspect:
        width = w
        height = w / aspect
    else:
        height = h
        width = h * aspect
    left = x + (w - width) / 2
    top = y + (h - height) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_header(slide, chapter, title, subtitle=None, accent=BLUE):
    add_tag(slide, chapter, 0.55, 0.32, color=accent)
    add_text(slide, title, 0.55, 0.74, 12.1, 0.52, 26, NAVY, True)
    add_rule(slide, 0.55, 1.34, 12.15, accent)
    if subtitle:
        add_text(slide, subtitle, 0.55, 1.48, 12.1, 0.28, 12, MUTED)


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_fade_transition(slide):
    transition = parse_xml(f'<p:transition {nsdecls("p")} spd="med"><p:fade/></p:transition>')
    slide._element.insert(2, transition)


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_fade_transition(slide)
    return slide


def title_slide(prs):
    s = blank(prs)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()
    add_text(s, "Financial AI Robustness Evaluation", 0.7, 0.82, 10.8, 0.78, 28, WHITE, True)
    add_text(
        s,
        "An integrated framework for downside-risk prediction, alert fusion, robustness, explainability, and cyber resilience",
        0.7,
        1.72,
        10.7,
        0.8,
        19,
        WHITE,
    )
    add_rule(s, 0.7, 2.72, 3.2, GOLD)
    add_text(s, "Dissertation defence presentation", 0.7, 3.02, 5.2, 0.32, 15, RGBColor(214, 223, 236), True)
    add_text(s, "Priyanka", 0.7, 5.82, 3.0, 0.3, 16, WHITE, True)
    add_text(s, "Chapter-wise defence deck", 0.7, 6.18, 4.2, 0.24, 12, RGBColor(214, 223, 236))
    add_text(s, "15-slide academic presentation", 9.55, 6.18, 3.0, 0.24, 12, RGBColor(214, 223, 236), align=PP_ALIGN.RIGHT)
    add_notes(
        s,
        "Open by stating the dissertation title and the central claim: this work is not only about predicting downside risk, but about evaluating whether a financial AI system remains useful when market conditions become disturbed. The claim supported by the full presentation is that an integrated pipeline reveals strengths and failure modes that a single headline metric would hide. Link forward by telling the examiner that the talk follows the dissertation chapter by chapter.",
    )


def roadmap_slide(prs):
    s = blank(prs)
    add_header(s, "Roadmap", "How the defence is organised", "The dissertation is presented in the same order as the written work.", GOLD)
    items = [
        ("1", "Problem", "Aim, gap, research questions"),
        ("2", "Method", "Data, features, models, evaluation design"),
        ("3", "Results", "Performance, alerts, robustness, explainability"),
        ("4", "Evaluation", "RQ answers, implications, recommendations"),
    ]
    xs = [0.75, 3.85, 6.95, 10.05]
    colors = [BLUE, TEAL, GOLD, RED]
    for (num, title, body), x, c in zip(items, xs, colors):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.55), Inches(2.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = c
        add_text(s, num, x + 0.16, 2.18, 0.35, 0.35, 22, c, True)
        add_text(s, title, x + 0.16, 2.66, 2.1, 0.3, 18, NAVY, True)
        add_text(s, body, x + 0.16, 3.16, 2.1, 0.76, 14, INK)
    add_text(s, "Chapter 4 receives the most time because it contains the evidence used to defend the dissertation.", 0.75, 5.2, 11.8, 0.3, 18, NAVY, True)
    add_notes(
        s,
        "Use this slide to control the room early. Explain that the structure mirrors the dissertation: first the problem, then the research design, then the evidence, then the evaluation and conclusion. The most important defence decision is that Chapter 4 receives the largest share of time because the examiner will care most about what the system actually achieved. Link forward by moving into the research problem.",
    )


def ch1_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 1", "Aim, problem, and research questions", "Why the study was needed before any modelling began.", BLUE)
    add_text(s, "Core problem", 0.65, 1.95, 2.0, 0.25, 17, NAVY, True)
    add_bullets(
        s,
        [
            "Financial AI systems can look accurate while missing rare but important downside events.",
            "Performance alone is insufficient if alerts, explanations, and robustness fail under disturbance.",
            "The dissertation asks whether these checks can be evaluated together in one framework.",
        ],
        0.65,
        2.28,
        5.6,
        2.0,
        18,
    )
    add_text(s, "Research questions", 6.75, 1.95, 2.2, 0.25, 17, NAVY, True)
    qbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(2.24), Inches(5.9), Inches(2.36))
    qbox.fill.solid()
    qbox.fill.fore_color.rgb = WHITE
    qbox.line.color.rgb = BLUE
    add_bullets(
        s,
        [
            "RQ1: Can AI predict downside-risk events?",
            "RQ2: Can monitoring and alerting identify suspicious behaviour?",
            "RQ3: Does the system remain useful under disturbed conditions?",
        ],
        6.96,
        2.48,
        5.45,
        1.72,
        17,
    )
    add_metric(s, "5 parts", "prediction → resilience", 0.75, 5.24, 1.55, BLUE)
    add_metric(s, "3 RQs", "directly evaluated", 2.55, 5.24, 1.55, TEAL)
    add_metric(s, "1", "shared pipeline", 4.35, 5.24, 1.55, GOLD)
    add_notes(
        s,
        "Start from the need, not from the code. The research problem is that financial AI can appear successful under ordinary metrics while still being unsafe or unhelpful when rare events matter most. State the three questions clearly because the rest of the defence answers them one by one. The slide supports the claim that the dissertation is evaluative, not merely a model-building exercise. Link forward by explaining what the literature had already solved and what it still left fragmented.",
    )


def ch2_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 2", "Literature review and research gap", "Existing studies solve pieces of the problem, but rarely evaluate them together.", TEAL)
    add_text(s, "What prior work usually covers", 0.65, 1.98, 3.3, 0.25, 17, NAVY, True)
    add_bullets(
        s,
        [
            "Prediction quality",
            "Explainability",
            "Adversarial robustness",
            "Drift monitoring",
            "Governance",
        ],
        0.65,
        2.46,
        2.8,
        2.0,
        18,
    )
    add_text(s, "Research gap", 3.55, 1.98, 2.0, 0.25, 17, NAVY, True)
    gap = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.55), Inches(2.25), Inches(3.15), Inches(2.35))
    gap.fill.solid()
    gap.fill.fore_color.rgb = WHITE
    gap.line.color.rgb = TEAL
    add_text(s, "The literature is strong in fragments,\nbut weaker on integrated evidence:\n\nCan the same system remain predictive,\nexplainable, alertable, and resilient\nunder disturbance?", 3.8, 2.53, 2.68, 1.7, 17, INK, True)
    add_picture_fit(s, CH4 / "figure_4_24_integrated_evaluation_matrix.png", 7.0, 1.95, 5.75, 3.25)
    add_text(s, "Used here as the conceptual bridge from literature gap to dissertation evaluation design.", 7.03, 5.24, 5.6, 0.3, 12, MUTED)
    add_notes(
        s,
        "This replaces the earlier mistake of showing a methodology pipeline in the literature-review section. Explain that Chapter 2 finds many useful strands in prior work, but the gap is their separation. The matrix on the right is used as a conceptual bridge: each research question maps to an evaluation layer and to evidence later produced in Chapter 4. A likely challenge is whether this is really novel; answer that the contribution is not a brand-new classifier, but the integrated evaluation of prediction, alerts, robustness, explanations, and governance-relevant artefacts. Link forward by showing how the methodology was built to close that gap.",
    )


def ch3_overview_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 3", "Methodology overview", "One integrated workflow rather than isolated experiments.", BLUE)
    add_picture_fit(s, CH3 / "figure_3_1_integrated_methodology_pipeline.png", 0.65, 1.95, 12.0, 4.95)
    add_notes(
        s,
        "This is the correct home for Figure 3.1. Walk through the pipeline from market data to feature engineering, model training, alert generation, robustness testing, explainability, and governance artefacts. Emphasise that each later result is produced from the same pipeline, which prevents disconnected claims. The slide supports the methodological integrity of the study. Link forward by zooming into the data, target, and chronological split.",
    )


def ch3_data_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 3", "Data, features, target, and chronological split", "The system is trained and tested in time order to avoid future information leakage.", TEAL)
    add_picture_fit(s, CH3 / "figure_3_2_preprocessing_feature_engineering_workflow.png", 0.55, 1.95, 7.4, 3.75)
    add_picture_fit(s, CH4 / "figure_4_3_dataset_split_target_distribution.png", 8.08, 1.95, 4.7, 3.75)
    add_bullets(
        s,
        [
            "Public market data transformed into technical, volatility, volume, and drawdown features.",
            "Target: financially meaningful future downside-risk event.",
            "Chronological split preserves the real deployment direction: past → future.",
        ],
        0.7,
        5.78,
        12.0,
        0.95,
        17,
    )
    add_notes(
        s,
        "Explain this in beginner-friendly language if asked: the system learns from earlier dates and is tested on later dates, just as it would be in practice. That is what chronological split means, and it avoids accidentally letting the model learn from the future. The left visual shows how raw market records become engineered features; the right visual shows the target distribution and split used in the experiment. Link forward by explaining the model families and the evaluation design.",
    )


def ch3_models_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 3", "Model development and evaluation design", "The study compares detectors, then deliberately perturbs the system.", GOLD)
    add_picture_fit(s, CH3 / "figure_3_3_detection_mechanism_comparison.png", 0.55, 1.92, 6.0, 3.9)
    add_picture_fit(s, CH3 / "figure_3_4a_six_part_robustness_evaluation_design.png", 6.78, 1.92, 5.95, 3.9)
    add_text(s, "Evaluation layers", 0.65, 5.95, 1.8, 0.25, 16, NAVY, True)
    add_bullets(
        s,
        [
            "Baseline classification",
            "Threshold optimisation",
            "Alert fusion",
            "Noise, stress, adversarial, explainability, drift, cyber-resilience checks",
        ],
        2.1,
        5.88,
        10.4,
        0.62,
        16,
    )
    add_notes(
        s,
        "This slide explains how the dissertation moves beyond one model and one score. The left figure distinguishes supervised and event-detection mechanisms; the right figure shows the six-part robustness evaluation. If challenged on why so many tests were needed, answer that a financial AI system can fail in more than one way: it may misclassify, degrade under noise, become unstable in its explanations, or trigger poor operational behaviour. Link forward into the results chapter, where these tests begin with the baseline models.",
    )


def ch4_baseline_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 4", "Baseline results and threshold interpretation", "Accuracy alone hid the real behaviour of the models.", RED)
    add_picture_fit(s, FIG / "supervised_model_test_metrics.png", 0.55, 1.88, 5.55, 3.35)
    add_picture_fit(s, FIG / "confusion_matrix_heatmaps_supervised.png", 6.18, 1.88, 6.55, 3.35)
    add_metric(s, "0.745", "LogReg recall", 0.72, 5.55, 1.45, BLUE)
    add_metric(s, "0.006", "RF recall", 2.38, 5.55, 1.45, RED)
    add_metric(s, "764 / 1025", "events found", 4.04, 5.55, 1.55, TEAL)
    add_text(s, "Interpretation: the seemingly stronger accuracy of Random Forest was misleading because it almost never found the rare downside events.", 6.25, 5.55, 6.1, 0.72, 17, NAVY, True)
    add_notes(
        s,
        "Lead with the main lesson: for imbalanced downside-risk classification, accuracy is the wrong headline metric. Logistic Regression found 764 of 1,025 test events with recall 0.745, while Random Forest achieved high accuracy but only 0.006 recall, meaning it nearly ignored the events that mattered most. This slide supports RQ1 and justifies the later threshold and alert analysis. A likely challenge is whether the weaker model should be discarded; answer that the result demonstrates why evaluation must be aligned to operational risk, not cosmetic accuracy. Link forward by showing how alert fusion turns model outputs into visible decisions.",
    )


def ch4_alert_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 4", "Event detection, alert fusion, and visible alert evidence", "The system produces concrete flags, not only aggregate metrics.", BLUE)
    add_picture_fit(s, CH4 / "figure_4_8_operational_alert_workflow.png", 0.55, 1.88, 5.75, 2.18)
    add_picture_fit(s, FIG / "alert_flags_nvda_example.png", 6.45, 1.88, 6.25, 3.58)
    add_metric(s, "0.957", "Fusion OR recall", 0.72, 4.42, 1.55, TEAL)
    add_metric(s, "0.357", "Consensus F1", 2.52, 4.42, 1.55, GOLD)
    add_text(
        s,
        "Visible example: flagged dates are overlaid on the stock path, making the alerting subsystem explainable in a presentation and auditable in practice.",
        0.72,
        5.56,
        5.45,
        0.8,
        17,
        NAVY,
        True,
    )
    add_notes(
        s,
        "This is the slide that makes the alert system tangible. Explain that alerting is not based solely on the five-day fall label; the system combines model outputs, anomaly evidence, and alert logic. The left workflow shows the operational process, and the right plot demonstrates real flagged instances on a stock chart. Fusion OR reached 0.957 recall, while consensus produced a more selective F1 of 0.357, so the trade-off is coverage versus alert burden. Link forward by asking the next defence question: what happens when the environment is deliberately disturbed?",
    )


def ch4_robustness_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 4", "Robustness, stress, and adversarial testing", "The study measures how performance changes when conditions worsen.", GOLD)
    add_picture_fit(s, FIG / "supervised_f1_drop_under_noise.png", 0.55, 1.92, 4.05, 3.55)
    add_picture_fit(s, FIG / "supervised_stress_scenario_f1_comparison.png", 4.72, 1.92, 4.05, 3.55)
    add_picture_fit(s, CH3 / "figure_3_4b_noise_stress_adversarial_defence_workflow.png", 8.88, 1.92, 3.85, 3.55)
    add_bullets(
        s,
        [
            "Noise and stress expose performance degradation.",
            "Adversarial tests reveal vulnerability beyond ordinary market variation.",
            "The contribution is measured resilience, not an assumption of resilience.",
        ],
        0.72,
        5.72,
        12.0,
        0.78,
        17,
    )
    add_notes(
        s,
        "Use this slide to show that robustness was actively tested rather than claimed. The first two figures show F1 degradation under noisy and stressed conditions; the workflow figure explains how the adversarial and defence logic fits into the evaluation. This supports RQ3 by demonstrating what changes when data quality or market behaviour worsens. If challenged on whether the model is robust, answer precisely: it is only partially robust, and the dissertation's value is that it quantifies the degradation instead of hiding it. Link forward by moving from output stability to explanation stability and monitoring over time.",
    )


def ch4_xai_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 4", "Explainability, time-to-detection, and drift monitoring", "Operational usefulness requires more than a prediction at one timestamp.", TEAL)
    add_picture_fit(s, CH4 / "figure_4_18_shap_stability.png", 0.55, 1.92, 4.2, 3.42)
    add_picture_fit(s, CH3 / "figure_3_4d_telemetry_time_to_detection_measurement.png", 4.92, 1.92, 3.75, 3.42)
    add_picture_fit(s, CH3 / "figure_3_4_robustness_adversarial_drift_framework.png", 8.84, 1.92, 3.9, 3.42)
    add_bullets(
        s,
        [
            "SHAP checks whether explanations remain stable under disturbance.",
            "Time-to-detection checks whether alerts arrive early enough to matter.",
            "Drift monitoring creates a governance trigger for review or retraining.",
        ],
        0.72,
        5.58,
        12.0,
        0.9,
        17,
    )
    add_notes(
        s,
        "Explain that a technically correct prediction is still not enough if the explanation becomes unreliable, the alert arrives too late, or the model quietly drifts after deployment. These three evaluations extend the dissertation from model performance into operational usefulness. The slide supports the broader RQ3 argument that robustness includes explanation stability and monitoring, not only classification scores. Link forward by showing how the dissertation also considered hostile cyber scenarios and implementation coverage.",
    )


def ch4_cyber_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 4", "Cybersecurity resilience and implementation coverage", "The prototype was evaluated as a system, not only as a classifier.", RED)
    add_picture_fit(s, CH3 / "figure_3_4c_cybersecurity_resilience_scorecard_layers.png", 0.55, 1.92, 4.35, 3.56)
    add_picture_fit(s, FIG / "cyber_attack_detection_rates_by_type.png", 5.04, 1.92, 3.65, 3.56)
    add_picture_fit(s, CH4 / "figure_4_24_integrated_evaluation_matrix.png", 8.86, 1.92, 3.86, 3.56)
    add_bullets(
        s,
        [
            "Cyber scenarios test suspicious behaviour coverage.",
            "Layered resilience separates what is implemented from what remains governance work.",
            "The integrated matrix confirms that every RQ is supported by explicit evidence.",
        ],
        0.72,
        5.72,
        12.0,
        0.8,
        17,
    )
    add_notes(
        s,
        "This slide protects you from an examiner treating the dissertation as only a stock-prediction exercise. It shows that cyber-resilience scenarios and implementation coverage were explicitly considered. The important distinction is that the work delivers a functional research prototype and deployment-ready artefacts, not a production-authorised trading system. A likely challenge is overclaiming readiness; answer that the dissertation itself states further engineering, validation, and governance are required. Link forward by summarising the findings against the three research questions.",
    )


def ch5_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 5", "RQ-by-RQ evaluation and broader implications", "The findings answer the questions, but with important trade-offs.", BLUE)
    rows = [
        ("RQ1", "Yes, but metric choice matters", "Useful downside-risk detection was achieved; accuracy alone misled."),
        ("RQ2", "Yes, with alert trade-offs", "Fusion improved coverage, while precision and alert burden still matter."),
        ("RQ3", "Partly, and measurably", "Robustness, explanations, and drift degraded under disturbance."),
    ]
    y = 1.95
    colors = [BLUE, TEAL, GOLD]
    for (rq, verdict, evidence), c in zip(rows, colors):
        panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(y), Inches(12.05), Inches(1.12))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = c
        add_text(s, rq, 0.9, y + 0.19, 0.7, 0.28, 20, c, True)
        add_text(s, verdict, 1.75, y + 0.16, 3.4, 0.28, 18, NAVY, True)
        add_text(s, evidence, 5.35, y + 0.18, 6.95, 0.46, 16, INK)
        y += 1.34
    add_text(s, "Broader implication: financial AI should be judged by prediction, monitoring, resilience, explanation stability, and governance together.", 0.72, 6.18, 11.8, 0.4, 18, NAVY, True)
    add_notes(
        s,
        "This slide is your synthesis slide. Answer the research questions directly rather than leaving the examiner to infer the answer. RQ1 is supported, but only when the right metrics are used. RQ2 is supported because alert fusion improves suspicious-behaviour coverage, though the precision burden remains. RQ3 is only partly positive because the system remains evaluable under disturbance, but the tests also reveal degradation. The dissertation's broader implication is that financial AI evaluation must be multidimensional. Link forward into the conclusion, contribution, limitations, and recommendations.",
    )


def ch6_slide(prs):
    s = blank(prs)
    add_header(s, "Chapter 6", "Conclusion, contribution, limitations, recommendations", "A strong conclusion is precise about both achievement and boundary.", TEAL)
    cols = [
        ("Contribution", ["Integrated evaluation framework", "Visible alert evidence", "Deployment-oriented artefacts"], BLUE),
        ("Limitations", ["US large-cap scope", "Public data source", "Prototype, not production"], RED),
        ("Recommendations", ["Use recall-sensitive metrics", "Monitor drift continuously", "Retain human review"], GOLD),
    ]
    xs = [0.65, 4.43, 8.21]
    for (title, bullets, c), x in zip(cols, xs):
        panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.45), Inches(3.55))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = c
        add_text(s, title, x + 0.18, 2.22, 2.9, 0.25, 18, c, True)
        add_bullets(s, bullets, x + 0.18, 2.72, 2.95, 1.75, 16)
    add_text(
        s,
        "Final claim: the dissertation demonstrates a usable research prototype that exposes when a financial AI system works, when it degrades, and what must be monitored before deployment.",
        0.72,
        5.92,
        11.95,
        0.72,
        18,
        NAVY,
        True,
    )
    add_notes(
        s,
        "Conclude with balance. State the contribution clearly, but also state the boundaries clearly so the examiner does not need to force them out of you. The work contributes an integrated evaluation framework and concrete artefacts, but it is limited by data scope and by prototype status. The recommendations follow directly from the evidence: use recall-sensitive metrics, maintain drift monitoring, and preserve human review. Link forward to the final closing slide.",
    )


def closing_slide(prs):
    s = blank(prs)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()
    add_text(s, "Thank you", 0.75, 1.16, 4.0, 0.6, 30, WHITE, True)
    add_text(s, "Questions and discussion", 0.75, 1.98, 5.2, 0.42, 21, RGBColor(221, 229, 239))
    add_rule(s, 0.75, 2.72, 2.2, GOLD)
    add_text(
        s,
        "Defence anchor:\nA financial AI system should be judged by what it predicts,\nwhat it flags,\nhow it degrades,\nand whether humans can still trust it.",
        0.75,
        3.22,
        8.0,
        1.72,
        22,
        WHITE,
        True,
    )
    add_notes(
        s,
        "Close cleanly. If the examiner opens with a broad question, return to the defence anchor on the slide: prediction, alerting, degradation, and trust. This final slide does not introduce new evidence; it gives the room one concise sentence to remember.",
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide(prs)
    roadmap_slide(prs)
    ch1_slide(prs)
    ch2_slide(prs)
    ch3_overview_slide(prs)
    ch3_data_slide(prs)
    ch3_models_slide(prs)
    ch4_baseline_slide(prs)
    ch4_alert_slide(prs)
    ch4_robustness_slide(prs)
    ch4_xai_slide(prs)
    ch4_cyber_slide(prs)
    ch5_slide(prs)
    ch6_slide(prs)
    closing_slide(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
